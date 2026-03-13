"""
Document metadata & content email harvesting module.

Extracts emails from:
  - PDF files (text + XMP/DocInfo metadata)
  - DOCX/ODT word processor documents
  - XLSX/ODS spreadsheets
  - PPTX presentations
  - Plain text / CSV files
  - HTML files (raw source)
  - EML / RFC 2822 email files

Also handles remote documents discovered by the crawler.
"""

import io
import logging
import os
from typing import Dict, Optional, Set

from emailharvest.utils.extractor import extract_emails_raw, _is_valid

logger = logging.getLogger(__name__)

# Optional heavy imports — gracefully degrade if not installed
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from pdfminer.high_level import extract_text as pdfminer_extract
    PDFMINER_AVAILABLE = True
except ImportError:
    PDFMINER_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class DocumentParser:
    def __init__(self):
        self.found_emails: Dict[str, Set[str]] = {}

    def parse_file(self, path: str) -> Dict[str, Set[str]]:
        """Auto-detect and parse a local file."""
        ext = os.path.splitext(path)[1].lower()
        source = f"file:{path}"

        try:
            with open(path, "rb") as f:
                data = f.read()
        except OSError as e:
            logger.warning("Cannot read %s: %s", path, e)
            return {}

        return self.parse_bytes(data, ext, source)

    def parse_bytes(self, data: bytes, ext: str, source: str) -> Dict[str, Set[str]]:
        """Parse raw bytes given a file extension hint."""
        if ext == ".pdf":
            self._parse_pdf(data, source)
        elif ext in (".docx", ".doc"):
            self._parse_docx(data, source)
        elif ext in (".xlsx", ".xls"):
            self._parse_xlsx(data, source)
        elif ext == ".pptx":
            self._parse_pptx(data, source)
        elif ext in (".txt", ".csv", ".md", ".rst", ".log"):
            self._parse_text(data.decode(errors="replace"), source)
        elif ext in (".html", ".htm"):
            from emailharvest.utils.extractor import extract_emails_from_html
            for email in extract_emails_from_html(data.decode(errors="replace")):
                self._add(email, source)
        elif ext == ".eml":
            self._parse_eml(data, source)
        else:
            # Heuristic: try UTF-8 text extraction
            try:
                text = data.decode("utf-8", errors="replace")
                self._parse_text(text, source)
            except Exception:
                pass

        return self.found_emails

    # ------------------------------------------------------------------

    def _add(self, email: str, source: str) -> None:
        email = email.lower().strip(".")
        if _is_valid(email):
            if email not in self.found_emails:
                self.found_emails[email] = set()
                logger.info("[DocParser] Found: %s (in %s)", email, source)
            self.found_emails[email].add(source)

    def _parse_text(self, text: str, source: str) -> None:
        for email in extract_emails_raw(text):
            self._add(email, source)

    def _parse_pdf(self, data: bytes, source: str) -> None:
        if PYMUPDF_AVAILABLE:
            try:
                doc = fitz.open(stream=data, filetype="pdf")
                # Metadata
                meta = doc.metadata or {}
                for v in meta.values():
                    if v:
                        for email in extract_emails_raw(str(v)):
                            self._add(email, source + ":metadata")
                # Text content
                full_text = []
                for page in doc:
                    full_text.append(page.get_text())
                self._parse_text("\n".join(full_text), source + ":content")
                doc.close()
                return
            except Exception as e:
                logger.debug("PyMuPDF failed for %s: %s", source, e)

        if PDFMINER_AVAILABLE:
            try:
                text = pdfminer_extract(io.BytesIO(data))
                self._parse_text(text or "", source + ":pdfminer")
            except Exception as e:
                logger.debug("pdfminer failed for %s: %s", source, e)

    def _parse_docx(self, data: bytes, source: str) -> None:
        if not DOCX_AVAILABLE:
            return
        try:
            doc = DocxDocument(io.BytesIO(data))
            # Core properties / metadata
            props = doc.core_properties
            for field in ["author", "last_modified_by", "description", "subject", "keywords", "comments"]:
                val = getattr(props, field, None)
                if val:
                    for email in extract_emails_raw(str(val)):
                        self._add(email, source + ":metadata")
            # Paragraph text
            text = "\n".join(p.text for p in doc.paragraphs)
            self._parse_text(text, source + ":content")
            # Tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for email in extract_emails_raw(cell.text):
                            self._add(email, source + ":table")
        except Exception as e:
            logger.debug("DOCX parse failed for %s: %s", source, e)

    def _parse_xlsx(self, data: bytes, source: str) -> None:
        if not OPENPYXL_AVAILABLE:
            return
        try:
            wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
            # Workbook properties
            if wb.properties:
                for field in ["creator", "lastModifiedBy", "description", "subject", "keywords"]:
                    val = getattr(wb.properties, field, None)
                    if val:
                        for email in extract_emails_raw(str(val)):
                            self._add(email, source + ":metadata")
            # Cell values
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                for row in ws.iter_rows(values_only=True):
                    for cell_val in row:
                        if cell_val and isinstance(cell_val, str):
                            for email in extract_emails_raw(cell_val):
                                self._add(email, source + f":{sheet_name}")
        except Exception as e:
            logger.debug("XLSX parse failed for %s: %s", source, e)

    def _parse_pptx(self, data: bytes, source: str) -> None:
        if not PPTX_AVAILABLE:
            return
        try:
            prs = Presentation(io.BytesIO(data))
            # Core properties
            props = prs.core_properties
            for field in ["author", "last_modified_by", "subject", "description", "keywords"]:
                val = getattr(props, field, None)
                if val:
                    for email in extract_emails_raw(str(val)):
                        self._add(email, source + ":metadata")
            # Slide text
            for i, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        for email in extract_emails_raw(shape.text):
                            self._add(email, source + f":slide{i+1}")
        except Exception as e:
            logger.debug("PPTX parse failed for %s: %s", source, e)

    def _parse_eml(self, data: bytes, source: str) -> None:
        """Parse RFC 2822 email files."""
        import email as email_lib
        try:
            msg = email_lib.message_from_bytes(data)
            # Headers with email addresses
            for header in ["From", "To", "Cc", "Bcc", "Reply-To", "Sender"]:
                val = msg.get(header, "")
                if val:
                    for em in extract_emails_raw(val):
                        self._add(em, source + f":header:{header}")
            # Body
            if msg.is_multipart():
                for part in msg.walk():
                    ct = part.get_content_type()
                    if ct in ("text/plain", "text/html"):
                        payload = part.get_payload(decode=True)
                        if payload:
                            self._parse_text(payload.decode(errors="replace"), source + ":body")
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    self._parse_text(payload.decode(errors="replace"), source + ":body")
        except Exception as e:
            logger.debug("EML parse failed for %s: %s", source, e)


def parse_remote_documents(
    urls: Set[str],
    session,
    delay: float = 1.0,
) -> Dict[str, Set[str]]:
    """
    Fetch and parse a set of remote document URLs.
    Returns combined email -> source mapping.
    """
    from emailharvest.utils.http import safe_get
    import os

    parser = DocumentParser()
    for url in urls:
        ext = os.path.splitext(url.split("?")[0])[1].lower()
        logger.info("[DocParser] Fetching remote doc: %s", url)
        resp = safe_get(session, url, delay=delay)
        if resp and resp.content:
            parser.parse_bytes(resp.content, ext, url)

    return parser.found_emails
